"""
NAS 스캔 결과 발표자료 생성 (혼합/범용 청중용)

산출물:
  - data/charts/06-rag-queue.png        (RAG 인덱싱 큐 분포 — 신규)
  - data/charts/07-context-band.png     (분류 맥락 분포 = 반복 업무 발자국 — 신규)
  - docs/04-operation/nas-스캔-발표-YYYYMMDD.pptx  (10장 발표자료)

기초자료: docs/04-operation/nas-스캔-카탈로그-20260526.md (카탈로그 v5)
원본 차트 스타일: 05-report.py make_charts() 와 동일 (matplotlib barh, Malgun Gothic)
"""
from __future__ import annotations
import sys
from datetime import datetime
from pathlib import Path

from common import connect_db, DATA_DIR

CHART_DIR = DATA_DIR / "charts"
CHART_DIR.mkdir(exist_ok=True)
PPTX_PATH = (
    Path(__file__).resolve().parent.parent.parent
    / "docs" / "04-operation"
    / f"nas-스캔-발표-{datetime.now().strftime('%Y%m%d')}.pptx"
)

# ── RAG 큐 분류 CASE (05-report.py classify_rag_queue 와 동일) ────────────
RAG_QUEUE_VIEW = """
CREATE OR REPLACE VIEW rag_queue AS
SELECT path, name, ext, top_dir, size, mtime, pdf_pages,
    CASE
        WHEN ext = 'pdf' AND pdf_has_text = TRUE THEN 'Q1_text_pdf'
        WHEN ext IN ('hwp','hwpx') AND hwp_parseable = TRUE THEN 'Q2_hwp_auto'
        WHEN ext IN ('docx','doc','xlsx','xls','pptx','ppt','txt','md','csv','rtf','tsv','html','htm','xml','json') THEN 'Q3_general_doc'
        WHEN ext = 'pdf' AND pdf_is_scan_likely = TRUE THEN 'Q4_ocr_pdf'
        WHEN ext IN ('hwp','hwpx') AND (hwp_parseable = FALSE OR hwp_version = '3.0') THEN 'Q5_hwp_manual'
        WHEN ext IN ('jpg','jpeg','png','tiff','tif','bmp','gif','webp') THEN 'Q6_image_ocr_optional'
        WHEN ext IN ('mp4','m4v','mov','avi','mkv','wmv','flv','mpg','mpeg') THEN 'X1_video'
        WHEN ext IN ('tib','vhdx','vmdk','iso','ova','dmg') THEN 'X2_backup_image'
        WHEN ext IN ('zip','rar','7z','tar','gz','bz2','egg') THEN 'X3_archive'
        WHEN ext IN ('dat','sdf','raw','bin','sgy','xtf','df047','df038','df037','df025','ruv','tuv','mat','hdf','h5','nc','kan','sil','sdl','wvp','wvs','wvr','oculus','pds','jsf','wcd') THEN 'X4_measurement_raw'
        WHEN ext IN ('exe','dll','msi','sys','bat','sh','cmd','jar') THEN 'X5_executable'
        WHEN name LIKE '~$%' OR ext IN ('tmp','bak','now','opt','wpa','blc') OR name IN ('Thumbs.db','.DS_Store','desktop.ini') THEN 'X6_temp_system'
        WHEN ext = 'log' THEN 'X7_log'
        WHEN ext IN ('wav','mp3','m4a','flac','ogg','aac','wma') THEN 'X8_audio'
        WHEN ext IN ('dwg','dxf','dwf','step','stp','iges','igs','prt','sldprt','sldasm','3dm') THEN 'X9_cad'
        WHEN (ext IS NULL OR ext = '') AND size >= 50000000 THEN 'X10_noext_large'
        ELSE 'Q7_unknown'
    END AS rag_queue
FROM files WHERE is_readable
"""


def q(con, sql, params=None):
    return con.execute(sql, params or []).fetchall()


# ════════════════════════════════════════════════════════════════════════
# Part A — 차트 2종 생성
# ════════════════════════════════════════════════════════════════════════
def make_new_charts(con):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams["font.family"] = "Malgun Gothic"
    plt.rcParams["axes.unicode_minus"] = False

    out = []

    # ── 차트 6: RAG 인덱싱 큐 분포 (파일 수, Q=인덱싱 / X=skip 색 구분) ──
    con.execute(RAG_QUEUE_VIEW)
    rows = q(con, """
        SELECT rag_queue, COUNT(*) AS cnt
        FROM rag_queue GROUP BY rag_queue ORDER BY cnt DESC
    """)
    labels_kr = {
        "Q1_text_pdf": "Q1 텍스트 PDF", "Q2_hwp_auto": "Q2 HWP 자동",
        "Q3_general_doc": "Q3 일반문서", "Q4_ocr_pdf": "Q4 스캔PDF(OCR)",
        "Q5_hwp_manual": "Q5 HWP 수동", "Q6_image_ocr_optional": "Q6 이미지",
        "Q7_unknown": "Q7 미분류", "X1_video": "X1 영상", "X2_backup_image": "X2 백업이미지",
        "X3_archive": "X3 압축", "X4_measurement_raw": "X4 계측raw", "X5_executable": "X5 실행파일",
        "X6_temp_system": "X6 임시/시스템", "X7_log": "X7 로그", "X8_audio": "X8 음성",
        "X9_cad": "X9 CAD/도면", "X10_noext_large": "X10 무확장자대용량",
    }
    names = [labels_kr.get(r[0], r[0]) for r in rows]
    cnts = [r[1] for r in rows]
    # Q(인덱싱 대상)=teal, X(skip)=gray
    colors = ["#1F7A8C" if n.startswith("Q") else "#B0BEC5" for n in [r[0] for r in rows]]

    fig, ax = plt.subplots(figsize=(11, 6.2))
    ax.barh(names[::-1], [c / 1e6 for c in cnts[::-1]], color=colors[::-1])
    ax.set_xlabel("파일 수 (백만 건)")
    ax.set_title("RAG 인덱싱 큐 자동 분류 — Q=AI가 읽을 수 있는 자료 / X=비대상", fontsize=13)
    for i, v in enumerate(cnts[::-1]):
        ax.text(v / 1e6, i, f" {v:,}", va="center", fontsize=8)
    # 범례
    from matplotlib.patches import Patch
    ax.legend(handles=[
        Patch(color="#1F7A8C", label="Q — 인덱싱 대상 (문서·텍스트)"),
        Patch(color="#B0BEC5", label="X — 비대상 (영상·계측raw·백업 등)"),
    ], loc="lower right", fontsize=9)
    p = CHART_DIR / "06-rag-queue.png"
    plt.tight_layout(); plt.savefig(p, dpi=110); plt.close()
    out.append(p)
    print(f"     ✓ {p.name}")

    # ── 차트 7: 분류 맥락 분포 (반복 업무 발자국) ──
    rows = q(con, """
        WITH g AS (
            SELECT hash_head1mb, size, COUNT(*) AS contexts
            FROM files WHERE hash_head1mb IS NOT NULL AND is_readable
            GROUP BY hash_head1mb, size
        )
        SELECT
            CASE
                WHEN contexts = 1 THEN '단일 위치\n(유일 자료)'
                WHEN contexts = 2 THEN '2개 맥락'
                WHEN contexts = 3 THEN '3개 맥락'
                WHEN contexts BETWEEN 4 AND 5 THEN '4~5개\n(반복 업무)'
                WHEN contexts BETWEEN 6 AND 10 THEN '6~10개\n(표준 첨부물)'
                ELSE '10개+\n(회사 표준양식)'
            END AS band,
            COUNT(*) AS asset_count,
            MIN(contexts) AS ord
        FROM g GROUP BY band ORDER BY ord
    """)
    bands = [r[0] for r in rows]
    counts = [r[1] for r in rows]
    # 단일=회색(유일 자료), 다맥락=점점 진한 teal/주황 (반복 업무 발자국 강조)
    bar_colors = ["#90A4AE", "#7FB3C4", "#4F97A8", "#1F7A8C", "#E8843A", "#C75B1E"]
    fig, ax = plt.subplots(figsize=(11, 5.8))
    bars = ax.bar(bands, counts, color=bar_colors[:len(bands)])
    ax.set_ylabel("자료 수 (unique)")
    ax.set_title("자료 분류 맥락 분포 — 반복되는 자료 = 반복 업무의 발자국", fontsize=13)
    for b, c in zip(bars, counts):
        ax.text(b.get_x() + b.get_width() / 2, c, f"{c:,}", ha="center", va="bottom", fontsize=9)
    ax.margins(y=0.15)
    # 주석
    ax.annotate("← 진짜 핵심 산출물·raw 데이터가\n   여기 묻혀 있을 가능성",
                xy=(0, counts[0]), xytext=(0.6, counts[0] * 0.78),
                fontsize=9, color="#555")
    ax.annotate("반복 업무 자동화 후보 →\n(선박서류·계약양식 등)",
                xy=(len(bands) - 1, counts[-1]), xytext=(len(bands) - 2.4, max(counts) * 0.45),
                fontsize=9, color="#C75B1E",
                arrowprops=dict(arrowstyle="->", color="#C75B1E"))
    p = CHART_DIR / "07-context-band.png"
    plt.tight_layout(); plt.savefig(p, dpi=110); plt.close()
    out.append(p)
    print(f"     ✓ {p.name}")

    return out


# ════════════════════════════════════════════════════════════════════════
# Part B — PPTX 생성
# ════════════════════════════════════════════════════════════════════════
def build_pptx(con):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    NAVY = RGBColor(0x0A, 0x2A, 0x43)
    TEAL = RGBColor(0x1F, 0x7A, 0x8C)
    ACCENT = RGBColor(0xE8, 0x84, 0x3A)
    GRAY = RGBColor(0x5A, 0x6B, 0x73)
    LIGHT = RGBColor(0xEE, 0xF3, 0xF6)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    FONT = "Malgun Gothic"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(BLANK)

    def textbox(slide, l, t, w, h, text, size=18, bold=False, color=NAVY,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
        return tb

    def rect(slide, l, t, w, h, fill, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
        sp.shadow.inherit = False
        return sp

    def content_header(slide, kicker, title):
        rect(slide, 0, 0, SW.inches, 1.15, NAVY)
        rect(slide, 0, 1.15, SW.inches, 0.06, ACCENT)
        textbox(slide, 0.55, 0.12, 11, 0.35, kicker, size=12, bold=True, color=ACCENT, font=FONT)
        textbox(slide, 0.5, 0.4, 12.3, 0.7, title, size=26, bold=True, color=WHITE, font=FONT,
                anchor=MSO_ANCHOR.MIDDLE)

    def footer(slide, page):
        textbox(slide, 0.5, 7.05, 8, 0.35, "오션테크 · NAS 디지털 자산 전수 스캔 결과 (2026-05-26)",
                size=9, color=GRAY, font=FONT)
        textbox(slide, 12.0, 7.05, 1.0, 0.35, str(page), size=10, color=GRAY,
                align=PP_ALIGN.RIGHT, font=FONT)

    def picture_fit(slide, img, l, t, max_w, max_h):
        """이미지를 비율 유지하며 영역 안에 중앙 배치."""
        from PIL import Image
        try:
            with Image.open(img) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = 1100, 620
        ratio = min(max_w / (iw / 96), max_h / (ih / 96))
        w = (iw / 96) * ratio
        h = (ih / 96) * ratio
        left = l + (max_w - w) / 2
        top = t + (max_h - h) / 2
        slide.shapes.add_picture(str(img), Inches(left), Inches(top), Inches(w), Inches(h))

    # ── 슬라이드 1: 표지 ──────────────────────────────────────────────
    s = add_slide()
    rect(s, 0, 0, SW.inches, SH.inches, NAVY)
    rect(s, 0, 4.55, SW.inches, 0.07, ACCENT)
    textbox(s, 0.9, 1.7, 11.5, 0.5, "오션테크(OceanTech) 디지털 자산 현황", size=18, bold=True, color=TEAL, font=FONT)
    textbox(s, 0.85, 2.3, 11.7, 1.7, "NAS 전수 스캔 결과", size=54, bold=True, color=WHITE, font=FONT)
    textbox(s, 0.9, 4.75, 11.5, 0.8,
            "회사 전체 디지털 자료를 빠짐없이 분석하고\nAI 지식 검색(RAG) 도입을 위한 청사진을 그렸습니다",
            size=18, color=RGBColor(0xC8, 0xD6, 0xDF), font=FONT, line_spacing=1.2)
    textbox(s, 0.9, 6.4, 11.5, 0.5, "2026년 5월 26일  ·  소스: NAS(192.168.0.220) 45.4TB",
            size=13, color=GRAY, font=FONT)

    # ── 슬라이드 2: 한눈에 보는 규모 ──────────────────────────────────
    s = add_slide()
    content_header(s, "SCALE  ·  규모", "한눈에 보는 규모")
    cards = [
        ("8,040,781", "총 파일 수", "약 804만 건"),
        ("30.4 TB", "총 용량", "NAS 45.4TB 중"),
        ("17종", "자동 분류 큐", "처리 경로별 정리"),
    ]
    cw, gap, x0, cy = 3.9, 0.35, 0.62, 2.1
    for i, (big, label, sub) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, cy, cw, 3.1, LIGHT)
        rect(s, x, cy, cw, 0.14, TEAL)
        textbox(s, x, cy + 0.55, cw, 1.2, big, size=44, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x, cy + 1.95, cw, 0.5, label, size=18, bold=True, color=TEAL,
                align=PP_ALIGN.CENTER)
        textbox(s, x, cy + 2.45, cw, 0.5, sub, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    textbox(s, 0.62, 5.6, 12.1, 1.0,
            "회사가 20여 년간 쌓아온 모든 디지털 자료를 한 건도 빠뜨리지 않고 스캔했습니다.\n"
            "단순 파일 목록이 아니라, 내용·중복·업무 맥락까지 분석한 \"자산 카탈로그\"입니다.",
            size=15, color=GRAY, font=FONT, line_spacing=1.3)
    footer(s, 2)

    # ── 슬라이드 3: 무엇이 들어있나 (확장자) ──────────────────────────
    s = add_slide()
    content_header(s, "WHAT  ·  자료 종류", "무엇이 들어있나 — 자료 종류별 용량")
    picture_fit(s, CHART_DIR / "01-asset-by-ext.png", 0.5, 1.45, 8.6, 5.3)
    rect(s, 9.35, 1.6, 3.5, 5.0, LIGHT)
    textbox(s, 9.6, 1.85, 3.0, 4.6,
            "읽어보면\n\n"
            "• 영상(mp4·avi) 이 용량 1위\n   — 현장·관측 영상\n\n"
            "• 사진(jpg) 258만 장\n   — 장비·현장 기록\n\n"
            "• 문서(txt·csv·hwp·pdf)\n   — AI가 읽을 핵심 자료\n\n"
            "• 계측 raw(dat) 130만\n   — 관측 원천 데이터",
            size=13, color=NAVY, font=FONT, line_spacing=1.15)
    footer(s, 3)

    # ── 슬라이드 4: 누가 무엇을 쌓았나 (부서·기관) ────────────────────
    s = add_slide()
    content_header(s, "WHO  ·  부서·기관", "누가 무엇을 쌓았나 — 부서·기관별 보유량")
    picture_fit(s, CHART_DIR / "02-asset-by-toplevel.png", 0.5, 1.45, 8.6, 5.3)
    rect(s, 9.35, 1.6, 3.5, 5.0, LIGHT)
    textbox(s, 9.6, 1.85, 3.0, 4.6,
            "핵심 고객·사업\n\n"
            "• 국립해양조사원(KHOA)\n   해양관측부이 유지관리\n   — 최대 자산군\n\n"
            "• 사업팀·기술팀\n   팀 단위 축적 자료\n\n"
            "• 기상청·수자원공사\n   해양과학기술원 등\n\n"
            "공공 해양관측 사업이\n자산의 중심축입니다.",
            size=13, color=NAVY, font=FONT, line_spacing=1.15)
    footer(s, 4)

    # ── 슬라이드 5: 언제 만들어졌나 (활동 패턴) ───────────────────────
    s = add_slide()
    content_header(s, "WHEN  ·  시기", "언제 만들어졌나 — 회사 활동 패턴")
    picture_fit(s, CHART_DIR / "03-activity-by-year.png", 0.5, 1.45, 8.0, 5.3)
    # Hot/Warm/Cold 요약
    rows = [
        ("최근 1년", "45만", "6,224 GB", "Hot — 우선 인덱싱", ACCENT),
        ("1~2년", "64만", "5,628 GB", "Warm", TEAL),
        ("2~5년", "230만", "9,441 GB", "Active archive", TEAL),
        ("5~10년", "261만", "8,288 GB", "Cold archive", GRAY),
        ("10년+", "204만", "3,860 GB", "보존 자산", GRAY),
    ]
    ty = 1.7
    textbox(s, 8.75, ty, 4.1, 0.4, "시기별 분포", size=15, bold=True, color=NAVY)
    ty += 0.55
    for period, cnt, gb, meaning, col in rows:
        rect(s, 8.75, ty, 0.12, 0.62, col)
        textbox(s, 9.0, ty, 1.5, 0.62, period, size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, 10.4, ty, 1.2, 0.62, cnt, size=13, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        textbox(s, 11.55, ty, 1.3, 0.62, meaning, size=10, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        ty += 0.7
    footer(s, 5)

    # ── 슬라이드 6: AI가 바로 읽을 수 있는 자료 ───────────────────────
    s = add_slide()
    content_header(s, "READY  ·  RAG 준비도", "AI가 바로 읽을 수 있는 자료")
    picture_fit(s, CHART_DIR / "06-rag-queue.png", 0.5, 1.45, 8.6, 5.3)
    rect(s, 9.35, 1.6, 3.5, 2.1, TEAL)
    textbox(s, 9.55, 1.78, 3.1, 0.5, "즉시 인덱싱 가능", size=15, bold=True, color=WHITE)
    textbox(s, 9.55, 2.25, 3.1, 1.0, "약 140만 파일", size=30, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 9.55, 3.15, 3.1, 0.5, "Q1 PDF + Q2 HWP + Q3 일반문서", size=11, color=RGBColor(0xD6,0xE7,0xEC))
    textbox(s, 9.35, 3.95, 3.5, 2.7,
            "• 자료를 처리 경로별로\n  자동 분류 (17종)\n\n"
            "• Q = AI가 읽을 수 있는\n  문서·텍스트\n\n"
            "• X = 영상·계측raw·백업\n  (검색 대상 아님,\n   메타데이터만)\n\n"
            "• 스캔 PDF 5.9만 건은\n  무료 OCR 후 추가",
            size=12, color=NAVY, font=FONT, line_spacing=1.1)
    footer(s, 6)

    # ── 슬라이드 7: 핵심 통찰 (반복 = 업무 발자국) ───────────────────
    s = add_slide()
    content_header(s, "★ INSIGHT  ·  핵심 통찰", "반복되는 자료 = 반복 업무의 발자국")
    picture_fit(s, CHART_DIR / "07-context-band.png", 0.5, 1.45, 8.4, 4.0)
    rect(s, 0.5, 5.55, 12.35, 1.35, LIGHT)
    textbox(s, 0.7, 5.65, 12.0, 1.2,
            "같은 자료가 여러 폴더에서 반복 등장한다는 것은 \"중요해서\"가 아니라 "
            "직원들이 같은 업무를 할 때마다 그 자료를 첨부·참조하기 때문입니다.\n"
            "→ 이 반복 패턴이 곧 ERP·지식플랫폼 업무 자동화 후보입니다. "
            "반대로 진짜 핵심 산출물(raw 데이터·최종보고서)은 '유일 자료' 32.6만 건에 묻혀 있습니다.",
            size=13, color=NAVY, font=FONT, line_spacing=1.25)
    footer(s, 7)

    # ── 슬라이드 8: 실제 사례 (선박 특별검사) ─────────────────────────
    s = add_slide()
    content_header(s, "★ CASE  ·  실제 사례", "실제 사례 — 선박 특별검사 서류 반복")
    textbox(s, 0.55, 1.4, 12.2, 0.95,
            "오션테크는 해양 조사 시 일반 어선을 임대해 작업선으로 용도 변경합니다(선박 특별검사 신청).\n"
            "그때마다 같은 선박서류를 매번 찾아 첨부 — 가장 뚜렷한 반복 업무 발자국입니다.",
            size=13, color=GRAY, font=FONT, line_spacing=1.2)
    # 표: 주력 임대선
    from pptx.util import Inches as IN
    rows_data = [
        ("선박명", "반복 등장", "사용 부서·기관", "성격"),
        ("만성호", "229 회", "13 곳", "★ 회사 주력 작업선"),
        ("일진호", "178 회", "11 곳", "주력 작업선"),
        ("씨로드호", "85 회", "12 곳", "주력 작업선"),
        ("은하호·스피드호·럭키2호·동진호 등", "기타", "다수", "프로젝트별 임대"),
    ]
    rcount, ccount = len(rows_data), 4
    tbl_shape = s.shapes.add_table(rcount, ccount, IN(0.55), IN(2.5), IN(7.6), IN(2.6))
    table = tbl_shape.table
    table.columns[0].width = IN(3.4)
    table.columns[1].width = IN(1.4)
    table.columns[2].width = IN(1.4)
    table.columns[3].width = IN(1.4)
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.name = FONT
            run.font.size = Pt(13 if ri == 0 else 12)
            run.font.bold = (ri == 0 or ri == 1)
            if ri == 0:
                run.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = NAVY
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    # 우측 자동화 박스
    rect(s, 8.5, 2.5, 4.35, 2.6, TEAL)
    textbox(s, 8.7, 2.65, 4.0, 0.5, "자동화 후보", size=15, bold=True, color=WHITE)
    textbox(s, 8.7, 3.15, 4.0, 1.9,
            "ERP에 '선박 자산 마스터'를\n만들고 선박별 표준 서류를\n묶어두면,\n\n"
            "프로젝트에서 \"작업선: 만성호\"\n선택 시 최신 선적증서·\n어선검사증서가 자동 첨부.\n"
            "검사증서 연차 갱신도 자동 관리.",
            size=12.5, color=WHITE, font=FONT, line_spacing=1.15)
    textbox(s, 0.55, 5.4, 12.2, 1.3,
            "이런 반복 업무 후보를 데이터로 발견했습니다 — 행정서류(선적·검사증서) 3,148건, "
            "양식 HWP 1,017건, 표준 사진 1.5만 장.\n"
            "직원의 반복 검색·첨부 시간을 줄이고 누락을 방지하는 가장 빠른 자동화 지점입니다.",
            size=13, color=GRAY, font=FONT, line_spacing=1.25)
    footer(s, 8)

    # ── 슬라이드 9: 이걸로 뭘 할 수 있나 ──────────────────────────────
    s = add_slide()
    content_header(s, "VALUE  ·  활용", "이걸로 무엇을 할 수 있나")
    items = [
        ("🔎", "전사 지식 검색(RAG)", "20년치 문서를 자연어로 질문 — \"○○사업 계약 조건?\" 즉시 응답"),
        ("📎", "ERP 자동 첨부", "발주·계약·납품 시 선박서류·양식을 자동 첨부 (반복 검색 시간 절감)"),
        ("📂", "표준 서류 중앙관리", "회사 표준 양식·증서를 한 곳에서 버전 통제"),
        ("🧭", "신입 온보딩", "\"이 업무엔 이 자료가 필요\" 자동 안내"),
        ("📈", "업무 흐름 역추적", "자료 분포로 어떤 업무가 어떤 자료를 요구하는지 발견 → 자동화 우선순위"),
    ]
    y = 1.65
    for icon, title, desc in items:
        rect(s, 0.6, y, 12.1, 0.92, LIGHT)
        textbox(s, 0.8, y, 0.9, 0.92, icon, size=26, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        textbox(s, 1.85, y + 0.08, 3.6, 0.78, title, size=16, bold=True, color=TEAL,
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, 5.5, y + 0.08, 7.0, 0.78, desc, size=13, color=NAVY,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += 1.04
    footer(s, 9)

    # ── 슬라이드 10: 다음 단계 ────────────────────────────────────────
    s = add_slide()
    content_header(s, "NEXT  ·  다음 단계", "다음 단계 — NAS 지식검색(RAG) 도입")
    steps = [
        ("1", "인덱싱 큐 검토", "즉시 처리 가능한 140만 문서 목록 확정 (Q1·Q2·Q3 CSV 준비 완료)"),
        ("2", "Hot tier 결정", "어느 기간·부서를 먼저 검색에 올릴지 선택 (최근 1년 우선 권장)"),
        ("3", "RAG 시범(PoC)", "별도 벡터DB + 무료 OCR로 1주 시범 → 검색 품질 평가"),
        ("4", "지식플랫폼 확장", "ERP 연동 + 반복업무 자동화 (선박·계약서 등) 단계 도입"),
    ]
    x0, cw, gap, cy = 0.62, 2.95, 0.18, 2.0
    for i, (num, title, desc) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, cy, cw, 3.4, LIGHT)
        rect(s, x, cy, cw, 0.9, TEAL)
        textbox(s, x, cy + 0.1, cw, 0.7, num, size=34, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + 0.15, cy + 1.05, cw - 0.3, 0.8, title, size=16, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)
        textbox(s, x + 0.2, cy + 1.85, cw - 0.4, 1.4, desc, size=12, color=GRAY,
                align=PP_ALIGN.CENTER, line_spacing=1.15)
    textbox(s, 0.62, 5.7, 12.1, 0.9,
            "스캔·분석은 완료됐습니다. 이제 \"무엇부터 검색에 올릴지\"만 정하면 시범 도입을 시작할 수 있습니다.",
            size=14, bold=True, color=NAVY, font=FONT, align=PP_ALIGN.CENTER)
    footer(s, 10)

    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))
    return PPTX_PATH


def main() -> int:
    con = connect_db()
    total = q(con, "SELECT COUNT(*) FROM files")[0][0]
    if total == 0:
        print("⚠ files 테이블 비어 있음.")
        return 1
    print(f"발표자료 생성 시작 (총 {total:,} 파일)\n")
    print("  1. 신규 차트 2종 생성...")
    make_new_charts(con)
    print("  2. PPTX 생성...")
    path = build_pptx(con)
    con.close()
    print(f"\n✓ 완료\n  → {path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
