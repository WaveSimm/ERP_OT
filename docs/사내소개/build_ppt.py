# 오션테크 ERP·OT 플랫폼 사내 소개 PPT 생성기 (python-pptx)
# 실행: python build_ppt.py  →  ERP-OT-플랫폼-소개.pptx
# 각 슬라이드에 "이 페이지의 목적"을 발표자 노트로 기록.

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 오션테크 테마 (해양: 네이비 + 틸) ──
NAVY   = RGBColor(0x0F, 0x2A, 0x43)
TEAL   = RGBColor(0x1E, 0x88, 0xA8)
AQUA   = RGBColor(0x2B, 0xB3, 0xC0)
LIGHT  = RGBColor(0xF1, 0xF6, 0xF9)
GRAY   = RGBColor(0x55, 0x66, 0x70)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1B, 0x2A, 0x35)
FONT   = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def _set(run, size, color, bold=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.name = FONT

def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h).text_frame

def rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def footer(slide, page):
    tf = box(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "오션테크 ERP·OT 통합 플랫폼  ·  사내 소개 자료 v1.0"
    _set(r, 9, GRAY)
    p2 = tf.add_paragraph()  # page no right
    # page number right-aligned in separate box
    tb = box(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    pp = tb.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run(); rr.text = str(page); _set(rr, 9, GRAY)

def content_slide(title, kicker, bullets, purpose, page):
    """일반 콘텐츠 슬라이드: 상단 제목바 + 불릿."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Inches(0.06), AQUA)
    # kicker
    tf = box(s, Inches(0.55), Inches(0.18), Inches(12), Inches(0.35))
    r = tf.paragraphs[0].add_run(); r.text = kicker; _set(r, 11, AQUA, True)
    # title
    tf = box(s, Inches(0.55), Inches(0.48), Inches(12.2), Inches(0.65))
    r = tf.paragraphs[0].add_run(); r.text = title; _set(r, 26, WHITE, True)
    # bullets
    tf = box(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.3))
    tf.word_wrap = True
    first = True
    for b in bullets:
        level, text = b if isinstance(b, tuple) else (0, b)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7); p.level = level
        if level == 0:
            rb = p.add_run(); rb.text = "▪ "; _set(rb, 15, TEAL, True)
            rt = p.add_run(); rt.text = text; _set(rt, 15, DARK, True)
        else:
            rt = p.add_run(); rt.text = "– " + text; _set(rt, 13, GRAY)
    footer(s, page)
    note(s, purpose)
    return s

# ───────────────────────── 1. 표지 ─────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.08), AQUA)
tf = box(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.5))
r = tf.paragraphs[0].add_run(); r.text = "오션테크 ERP·OT 통합 플랫폼"; _set(r, 44, WHITE, True)
tf = box(s, Inches(0.95), Inches(3.7), Inches(11.5), Inches(0.7))
r = tf.paragraphs[0].add_run(); r.text = "프로젝트 · 자원 · 수리 · 회계 · 결재 · 지식검색을 하나로"; _set(r, 19, AQUA)
tf = box(s, Inches(0.95), Inches(4.8), Inches(11.5), Inches(0.6))
r = tf.paragraphs[0].add_run(); r.text = "사내 소개 자료 — 전 임직원 대상"; _set(r, 14, RGBColor(0xB8,0xCF,0xDD))
tf = box(s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "2026-06  ·  v1.0"; _set(r, 12, RGBColor(0x8FA,0x0,0x0) if False else RGBColor(0x90,0xA8,0xB8))
note(s, "목적: 표지. 발표 주제와 대상(전 임직원)을 제시하고 자료의 성격(사내 소개)을 명확히 한다.")

# ───────────────────────── 2. 발표 구성 (페이지별 목적) ─────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, Inches(1.15), NAVY); rect(s, 0, Inches(1.15), SW, Inches(0.06), AQUA)
tf = box(s, Inches(0.55), Inches(0.18), Inches(12), Inches(0.35))
r = tf.paragraphs[0].add_run(); r.text = "AGENDA"; _set(r, 11, AQUA, True)
tf = box(s, Inches(0.55), Inches(0.48), Inches(12), Inches(0.65))
r = tf.paragraphs[0].add_run(); r.text = "발표 구성과 페이지별 목적"; _set(r, 26, WHITE, True)
agenda = [
    ("1", "한눈에 보기 / 배경·목적", "무엇이고 왜 만들었는지"),
    ("2", "도입 효과 / 전체 구성 지도", "기대 효과와 모듈 전경"),
    ("3", "기능 모듈 상세", "프로젝트·자원·장비·수리·회계·결재·게시판"),
    ("4", "OT-Brain NAS 통합검색", "수십 년 자료를 검색 자산으로"),
    ("5", "기술·인프라·보안", "어떻게 안전하게 운영되는가"),
    ("6", "운영 현황 / 로드맵 / 문의", "지금 상태와 앞으로"),
]
tf = box(s, Inches(0.8), Inches(1.55), Inches(11.8), Inches(5.2)); tf.word_wrap = True
first = True
for no, t, purpose in agenda:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    p.space_after = Pt(12)
    r = p.add_run(); r.text = f"  {no}.  "; _set(r, 17, AQUA, True)
    r = p.add_run(); r.text = t; _set(r, 17, DARK, True)
    r = p.add_run(); r.text = f"   —  {purpose}"; _set(r, 13, GRAY)
footer(s, 2)
note(s, "목적: 발표 전체 구성과 각 묶음의 목적을 미리 제시해 청중이 흐름을 예측하게 한다. 각 슬라이드 노트에 페이지별 목적을 기재함.")

# ───────────────────────── 3. 한눈에 보기 ─────────────────────────
content_slide(
    "한눈에 보기", "WHAT IS IT",
    [
        "해양 관측·계측 사업(부이·관측소·선박)을 위한 맞춤형 통합 업무 시스템",
        (1, "무엇: 프로젝트·자원·수리·회계·결재·지식검색을 하나의 웹에서"),
        (1, "누가: 오션테크 전 임직원 (부서·직급별 권한 차등)"),
        (1, "어디서: 사내 웹 브라우저 접속 — 설치 불필요"),
        (1, "특징: 해양 도메인 특화 · 외부 SaaS 미사용 · 데이터 사내 보관"),
    ],
    "목적: 시스템의 정체성을 한 문장으로 각인시킨다. 세부 기능 전에 '전체상'을 먼저 잡아준다.", 3)

# ───────────────────────── 4. 배경과 목적 ─────────────────────────
content_slide(
    "도입 배경과 목적", "WHY",
    [
        "기존 문제",
        (1, "프로젝트 일정·인력·장비가 개별 엑셀·담당자에 분산"),
        (1, "수리·발주·재고·정산이 단절 → 같은 정보 중복 입력"),
        (1, "수십 년 NAS 자료가 폴더에 묻혀 검색 불가"),
        "목표",
        (1, "단일 진실 원천 — 프로젝트·자원·자산을 한 곳에서"),
        (1, "업무 연결 — 프로젝트→자원→장비→수리→회계→결재"),
        (1, "지식 자산화 + 사내 자체 운영(데이터 주권)"),
    ],
    "목적: '왜 필요한가'를 공감시킨다. 현재 불편(분산·중복·검색불가)을 짚고 해결 목표로 연결한다.", 4)

# ───────────────────────── 5. 도입 효과 ─────────────────────────
content_slide(
    "도입 효과", "BENEFITS",
    [
        "프로젝트 — 일정·진척·자원 배정 실시간 가시화, 진도율 자동 집계",
        "자원·장비 — 가용성·과부하 한눈에(피크/과배정 경고)",
        "수리(AS) — 접수→수리→부품→정산 이력 일원화 (현재 실운영)",
        "회계·구매 — 고객사·발주·재고·정산 연결, 중복 입력 제거",
        "전자결재 — 결재선·이력 표준화, 웹 어디서나",
        "지식검색 — NAS 800만 파일을 한국어·영어로 즉시 검색",
        "데이터 주권 — 모든 데이터 사내 보관, 외부 유출 위험 최소화",
    ],
    "목적: 도입으로 얻는 구체적 이득을 영역별로 제시해 '나에게 무슨 도움인가'를 보여준다.", 5)

# ───────────────────────── 6. 전체 구성 지도 ─────────────────────────
content_slide(
    "전체 구성 지도", "BIG PICTURE",
    [
        "업무 모듈 (웹 메뉴 9개 + 관리자)",
        (1, "내 대시보드 · 지휘센터(전사 현황)"),
        (1, "프로젝트 · 자원 · 장비 · 수리(AS)"),
        (1, "회계(구매·재고·발주·정산) · 전자결재 · 게시판+NAS검색"),
        "지식 플랫폼 — OT-Brain (NAS 800만 파일 검색)",
        "공통 기반 — 인증/권한 · 통합검색 · 첨부 · 알림 · 감사로그",
        "기술 기반 — Next.js · Node 마이크로서비스 · PostgreSQL · Docker(온프레미스)",
    ],
    "목적: 개별 기능 설명 전, 시스템 전체 레이어(업무/지식/공통/기술)를 지도로 보여줘 위치감을 준다.", 6)

# ───────────────────────── 7. 프로젝트 관리 ─────────────────────────
content_slide(
    "기능 ① 프로젝트 관리", "MODULE",
    [
        "해양 사업의 일정·자원을 관리하는 핵심 모듈",
        (1, "계층: 프로젝트 → 작업(Task) → 세그먼트(실제 일정+자원)"),
        (1, "자원 배정: 세그먼트별 인력·장비를 %/시간 단위로"),
        (1, "진도율 자동 집계: 하위 진척 → 상위 프로젝트 반영"),
        (1, "템플릿: 반복 사업을 일정 오프셋으로 빠르게 생성"),
        (1, "그룹핑: 부서·고객사·커스텀 기준 묶음(N:M)"),
    ],
    "목적: 가장 핵심 모듈인 프로젝트 관리의 구조(프로젝트-작업-세그먼트)와 자동 집계 강점을 설명한다.", 7)

# ───────────────────────── 8. 자원 & 장비 ─────────────────────────
content_slide(
    "기능 ② 자원 · 장비 관리", "MODULE",
    [
        "자원(사람) 관리",
        (1, "인력 가용성·부하율 시각화, 피크/과배정 경고"),
        (1, "휴가·근태 연동으로 실제 가용 인력 반영"),
        "장비 관리",
        (1, "장비를 독립 자산 풀로 관리(센서·계측기)"),
        (1, "체크아웃/반납 순환 — 프로젝트 대여·회수"),
        (1, "장비 마스터: 모델명·제조사·SKU·재고 통일 표기"),
    ],
    "목적: 인력과 장비라는 두 자원의 운용 방식과 '과부하 가시화' 가치를 함께 설명한다.", 8)

# ───────────────────────── 9. 수리 AS ─────────────────────────
content_slide(
    "기능 ③ 수리(AS) 관리  ★ 실운영 중", "MODULE — IN PRODUCTION",
    [
        "회사에서 가장 먼저 정식 운영에 들어간 모듈",
        (1, "접수 → 수리 → 부품 사용 → 정산 이력 일원화"),
        (1, "고객사·공급사·부품 마스터 연계"),
        (1, "통계로 수리 추이·부품 소진 파악"),
        (1, "엑셀로 관리하던 기존 수리 데이터 이관 완료"),
    ],
    "목적: 이미 실운영 중인 모듈을 강조해 '검증된 시스템'이라는 신뢰를 준다. 성공 사례로 활용.", 9)

# ───────────────────────── 10. 회계 ─────────────────────────
content_slide(
    "기능 ④ 회계 (구매·재고·발주·정산)", "MODULE",
    [
        "구매-재고-결재를 연결하는 통합 회계 모듈 (하위 12개 기능)",
        (1, "마스터: 고객사 · 공급사 · 품목(SKU) · 창고/위치"),
        (1, "흐름: 발주 → 입고 → 재고 → 번들출하 → 정산"),
        (1, "발주 워크플로우: 결재라인·회계정산·발주송금·부분결제 자동분할"),
        (1, "경비 접수(재무팀 연계) · 계약 · 변경 감사"),
    ],
    "목적: 가장 하위 기능이 많은 회계 모듈을 '마스터-흐름' 구조로 압축해 복잡도를 정리해 보여준다.", 10)

# ───────────────────────── 11. 결재 & 게시판 ─────────────────────────
content_slide(
    "기능 ⑤ 전자결재 · 게시판 · 협업", "MODULE",
    [
        "전자결재",
        (1, "경비정산·지출결의서 등(EXPENSE 양식 일원화)"),
        (1, "부서·직급 기반 결재선, 상신→승인/반려→완료 이력"),
        "게시판 · 협업",
        (1, "카테고리 게시판 · 기능요구 게시판(개선 트래킹)"),
        (1, "작업 비고(협업 메모) · 회사 달력(공휴일 자동 갱신)"),
    ],
    "목적: 업무 결재와 사내 협업 채널을 묶어, 시스템이 '소통·승인'까지 포함함을 보여준다.", 11)

# ───────────────────────── 12. 근태 & 관리자 ─────────────────────────
content_slide(
    "기능 ⑥ 근태 · 관리자 기능", "MODULE",
    [
        "근태·인사",
        (1, "출퇴근·근무 일정 관리, 휴가·휴일근무 자동 동기화"),
        (1, "인력 현황·부하율 계산과 연계"),
        "관리자 기능",
        (1, "사용자·권한 · 결재라인 · 회사달력 · 장비자원"),
        (1, "기능요청 트래킹 · 활동(감사) 로그"),
    ],
    "목적: 일반 사용자 기능 외 운영·관리 기능을 보여줘 시스템이 '관리 체계'까지 갖췄음을 전달한다.", 12)

# ───────────────────────── 13. OT-Brain (1) ─────────────────────────
content_slide(
    "OT-Brain — NAS 통합 지식검색 ①", "KNOWLEDGE",
    [
        "수십 년 NAS 자료를 검색 가능한 지식 자산으로 (게시판 → 🔎 NAS 통합검색)",
        (1, "규모: 약 804만 파일 / 30.4TB 스캔 완료"),
        (1, "파일명·폴더는 전량 즉시 검색, 본문은 단계적 추출"),
        "검색 방식 (한국어·영어 모두)",
        (1, "키워드(부분일치+단어경계) + 의미검색(AI 다국어 임베딩)"),
        (1, "중요도 가중(IDF) + GPU 재순위로 정확도 향상"),
    ],
    "목적: 차별화 포인트인 NAS 통합검색의 규모와 검색 원리를 설명해 '묻힌 자료를 깨운다'는 가치를 전달.", 13)

# ───────────────────────── 14. OT-Brain (2) ─────────────────────────
content_slide(
    "OT-Brain — NAS 통합 지식검색 ②", "KNOWLEDGE",
    [
        "파일 열기 (신규)",
        (1, "검색결과 파일명 클릭 → PDF·이미지·텍스트 브라우저 미리보기"),
        (1, "엑셀·한글·도면 등은 다운로드 / 폴더 경로 복사 지원"),
        "운영 방식",
        (1, "비용 0 — 로컬 AI 사용, 외부 전송 없음(데이터 사내 보관)"),
        (1, "본문 추출은 야간 자동 처리(업무시간 검색 우선)"),
        (1, "향후 ERP와 양방향 연동되는 통합 지식 플랫폼으로 확장"),
    ],
    "목적: 검색에서 끝나지 않고 '바로 열어보는' 실사용 경험과 비용 0·사내보관 운영 강점을 전달한다.", 14)

# ───────────────────────── 15. 기술 & 아키텍처 ─────────────────────────
content_slide(
    "기술 스택과 아키텍처", "TECH",
    [
        "프론트: Next.js 14 · TypeScript · Tailwind",
        "백엔드: Node.js 20 · Fastify · Prisma · Zod (마이크로서비스)",
        "데이터: PostgreSQL 16(서비스별 스키마) · Redis · RabbitMQ",
        "AI 검색: Ollama(bge-m3·qwen) + pgvector + pg_trgm",
        "구조: 모노레포(Turborepo) · 클린 아키텍처 4계층",
        "품질: 문서 우선 + PDCA(계획-설계-실행-검증) 주기",
    ],
    "목적: 기술 청중·경영진에게 '검증된 현대적 스택과 체계적 품질관리'로 신뢰를 준다. 과도한 세부는 생략.", 15)

# ───────────────────────── 16. 인프라 ─────────────────────────
content_slide(
    "인프라 구성", "INFRA",
    [
        "약 16개 컨테이너가 사내 서버 1대(Docker)에서 운영",
        (1, "PostgreSQL · Redis · RabbitMQ + 8개 마이크로서비스"),
        (1, "AI 검색(Ollama) · OCR 3종 · DB 자동 백업"),
        "온프레미스 — 외부 클라우드 비용 없음, GPU 1대로 AI 처리",
        "자동 백업 + 수동 스냅샷 보관",
    ],
    "목적: 외부 비용 없이 사내 1대 서버로 운영되는 효율적·자립적 인프라임을 보여준다.", 16)

# ───────────────────────── 17. 보안 ─────────────────────────
content_slide(
    "보안과 데이터 보호", "SECURITY",
    [
        "사내 전용 운영 — 현재 외부 인터넷 미노출(로컬 네트워크 한정)",
        "인증: JWT(Access 1h / Refresh 7d), 비밀번호 해싱",
        "권한 분리: 역할·직급별 메뉴·기능 접근 제어",
        "데이터 주권: 업무 데이터·NAS 자료 사내 보관(외부 AI 미전송)",
        "감사 로그 + 보안 일괄 패치(Critical 0건 유지)",
    ],
    "목적: 보안·데이터 보호에 대한 우려를 선제적으로 해소한다. 특히 '데이터가 밖으로 나가지 않는다'를 강조.", 17)

# ───────────────────────── 18. 운영 현황 ─────────────────────────
content_slide(
    "현재 운영 현황과 성숙도", "STATUS",
    [
        "수리(AS) 관리: ✅ 정식 운영(production) — 실제 업무 데이터로 운용",
        "그 외 모듈: 사전운영(pre-prod) — 안정화 후 점진적 정식 전환",
        "접속: 사내 네트워크 브라우저(외부 미노출)",
        "안정 모듈 다수: 프로젝트·근태·대시보드·인증·게시판·검색·발주 등",
        "진행 중: 경비정산 통합 · 재고 SKU · OCR 문서인식",
    ],
    "목적: 현재 어디까지 왔는지 솔직히 공유한다. '검증된 부분(AS)'과 '확장 중인 부분'을 구분해 신뢰 형성.", 18)

# ───────────────────────── 19. 로드맵 ─────────────────────────
content_slide(
    "향후 로드맵", "ROADMAP",
    [
        "통합 지식 플랫폼 본격화 — NAS·ERP·미래 자료 통합",
        "NAS 검색 고도화 — 본문 추출 확대 · 한↔영 동의어 · 증분 자동 갱신",
        "모듈 정식 전환 — pre-prod → production 단계적 전환",
        "운영 최적화 — 백업 확장 · 오프사이트 · TLS 등 이관 준비",
        "AI 활용 확대 — 문서 질의응답(RAG) · OCR 자동 추출(GPU 증설 시)",
    ],
    "목적: 앞으로의 방향을 제시해 기대감을 주고, 사용자 피드백이 로드맵에 반영됨을 알린다.", 19)

# ───────────────────────── 20. 마무리 ─────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.4), SW, Inches(0.08), AQUA)
tf = box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2))
r = tf.paragraphs[0].add_run(); r.text = "함께 만들어가는 시스템입니다"; _set(r, 34, WHITE, True)
tf = box(s, Inches(0.95), Inches(3.5), Inches(11.5), Inches(1.6)); tf.word_wrap = True
for line in [
    "접속: 사내 네트워크 → ERP 로그인 → 메뉴 이용",
    "NAS 검색: 게시판 상단 🔎 NAS 통합검색",
    "문의·개선 요청: 게시판 → 기능 요청 또는 담당자",
]:
    p = tf.add_paragraph(); p.space_after = Pt(8)
    r = p.add_run(); r.text = "·  " + line; _set(r, 16, AQUA)
note(s, "목적: 마무리. 실제 접속 방법과 피드백 채널을 안내해 '바로 써보고 의견 달라'는 행동을 유도한다.")

prs.save("ERP-OT-플랫폼-소개.pptx")
print("saved: ERP-OT-플랫폼-소개.pptx  /  slides =", len(prs.slides._sldIdLst))
